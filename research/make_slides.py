"""
Builds the presentation deck: research/NutriCoach-T2D-slides.pptx

Edit the text below and re-run to regenerate (needs: pip install python-pptx).
Design: white content slides framed by dark-teal title/closing slides; the
traffic-light dots of the app as recurring motif; study-arm colors match the
figures (blue = control, orange = intervention). All numbers are the real
outputs of the pipeline (test suite, evaluation, analyze.py) - the study data
are simulated and labeled as such on every results slide.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "NutriCoach-T2D-slides.pptx")
FIG = os.path.join(HERE, "output")

TEAL = RGBColor(0x0E, 0x7C, 0x66)
TEAL_DARK = RGBColor(0x0A, 0x5C, 0x4C)
TEAL_PALE = RGBColor(0xCF, 0xE8, 0xE1)
INK = RGBColor(0x1C, 0x2B, 0x28)
MUTED = RGBColor(0x5F, 0x73, 0x6E)
TINT = RGBColor(0xEE, 0xF5, 0xF3)
LINE = RGBColor(0xDB, 0xE6, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1A, 0x98, 0x50)
YELLOW = RGBColor(0xD9, 0xA4, 0x00)
RED = RGBColor(0xD6, 0x45, 0x41)
ORANGE = RGBColor(0xE6, 0x9F, 0x00)
ORANGE_DARK = RGBColor(0xA6, 0x6B, 0x00)
ORANGE_TINT = RGBColor(0xFD, 0xF3, 0xE1)
BLUE = RGBColor(0x00, 0x72, 0xB2)
BLUE_TINT = RGBColor(0xE3, 0xEF, 0xF7)
RED_TINT = RGBColor(0xFB, 0xEA, 0xE9)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]


def slide_new(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def tb(s, x, y, w, h, paras, wrap=True):
    """Text box from a list of paragraph dicts: text, size, bold, color, font, italic, align, after."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        para.space_after = Pt(p.get("after", 4))
        run = para.add_run()
        run.text = p["text"]
        f = run.font
        f.size = Pt(p.get("size", 13))
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = p.get("color", INK)
        f.name = p.get("font", "Calibri")
    return box


def shape(s, kind, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and sp.adjustments:
        sp.adjustments[0] = radius
    sp.text_frame.word_wrap = True
    return sp


def dots(s, x, y, r=0.09, gap=0.28):
    for i, col in enumerate((GREEN, YELLOW, RED)):
        shape(s, MSO_SHAPE.OVAL, x + i * gap, y, r, r, fill=col)


def title(s, text, width=9.0):
    tb(s, 0.5, 0.32, width, 0.62, [{"text": text, "size": 30, "bold": True, "font": "Cambria"}])
    dots(s, 8.93, 0.42)


def footer(s, num, dark=False):
    col = RGBColor(0x8F, 0xBF, 0xB2) if dark else MUTED
    tb(s, 0.5, 5.30, 6.0, 0.24, [{"text": "NutriCoach T2D · proof of concept · simulated data", "size": 9, "color": col}])
    tb(s, 9.0, 5.30, 0.5, 0.24, [{"text": str(num), "size": 9, "color": col, "align": PP_ALIGN.RIGHT}])


def badge_simulated(s):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.45, 0.40, 3.05, 0.36, fill=ORANGE_TINT, line=ORANGE, radius=0.5)
    tb(s, 6.45, 0.465, 3.05, 0.26, [{"text": "SIMULATED DATA - METHODOLOGY DEMO", "size": 9.5, "bold": True,
                                     "color": ORANGE_DARK, "align": PP_ALIGN.CENTER}])


def bullets(items, size=12.5, color=INK, after=7):
    return [{"text": "•  " + t, "size": size, "color": color, "after": after} for t in items]


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- 1 · title
s = slide_new(TEAL_DARK)
dots(s, 0.55, 1.22, r=0.14, gap=0.40)
tb(s, 0.5, 1.55, 8.6, 0.95, [{"text": "NutriCoach T2D", "size": 44, "bold": True, "color": WHITE, "font": "Cambria"}])
tb(s, 0.5, 2.55, 7.2, 1.0, [{"text": "A rule-based conversational AI for nutrition self-management support in type 2 diabetes",
                             "size": 17, "color": TEAL_PALE, "after": 8},
                            {"text": "Proof of concept with a simulated evaluation pipeline", "size": 12.5,
                             "italic": True, "color": TEAL_PALE}])
tb(s, 0.5, 4.62, 7.0, 0.4, [{"text": "Sana Sajid · Research project", "size": 13, "bold": True, "color": WHITE}])
for i, col in enumerate((GREEN, YELLOW, RED)):
    shape(s, MSO_SHAPE.OVAL, 8.55, 1.75 + i * 0.85, 0.52, 0.52, fill=col)
notes(s, "One sentence to open: I built a working chatbot that supports people with type 2 diabetes in everyday "
         "nutrition decisions - and a complete statistical pipeline showing how its effect would be evaluated. "
         "Everything runs locally, and all study data are simulated and labeled as such.")

# ---------------------------------------------------------------- 2 · problem
s = slide_new()
title(s, "Why nutrition support needs to scale")
tb(s, 0.5, 1.35, 4.2, 1.6, [{"text": "537 M", "size": 60, "bold": True, "color": TEAL, "font": "Cambria"},
                            {"text": "adults were living with diabetes in 2021 (IDF Diabetes Atlas) - roughly 9 in 10 have type 2.",
                             "size": 12, "color": MUTED}])
tb(s, 0.5, 3.35, 4.2, 1.7, [{"text": "Nutrition is first-line therapy - but the questions arrive daily, between appointments.",
                             "size": 14, "bold": True, "color": INK}])
card = shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.1, 1.35, 4.4, 3.6, fill=TINT, line=LINE, radius=0.06)
tb(s, 5.45, 1.65, 3.75, 3.1, [{"text": "The daily reality", "size": 15, "bold": True, "after": 10}] + bullets([
    "Every meal is a decision: \"Can I eat this? How much? What instead?\"",
    "Dietitian access is limited and episodic",
    "Generic leaflets don't answer the question of the moment",
    "A chatbot answers 24/7, at near-zero marginal cost",
], size=12.5, after=9))
footer(s, 2)
notes(s, "Frame the gap: nutrition therapy is first-line in every guideline, but professional support is scarce "
         "and episodic while the decisions are daily. That mismatch is what conversational agents can close.")

# ---------------------------------------------------------------- 3 · why rule-based
s = slide_new()
title(s, "Design decision: why not just an LLM?")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 1.3, 4.4, 3.3, fill=TINT, line=LINE, radius=0.06)
tb(s, 0.85, 1.6, 3.75, 2.9, [{"text": "LLM chatbot", "size": 15, "bold": True, "color": MUTED, "after": 10}] + bullets([
    "Can hallucinate confident nonsense",
    "Answers vary from run to run",
    "Safety lives in a prompt, not in code",
    "Input usually leaves the device",
], color=INK, after=9))
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.1, 1.3, 4.4, 3.3, fill=TEAL, radius=0.06)
tb(s, 5.45, 1.6, 3.75, 2.9, [{"text": "Rule-based + curated knowledge", "size": 15, "bold": True, "color": WHITE, "after": 10}] + bullets([
    "Every answer grounded in reviewed content",
    "Deterministic and reproducible",
    "Safety guaranteed by code, checked first",
    "100% offline - privacy by default",
], color=WHITE, after=9))
tb(s, 0.5, 4.78, 9.0, 0.45, [{"text": "Trade-off: limited coverage - measured via the fallback rate. A hybrid (LLM language + curated content) is the designed next step.",
                              "size": 11.5, "italic": True, "color": MUTED}])
footer(s, 3)
notes(s, "This is the research angle, not a limitation of skill: in a health context, reproducibility, auditability "
         "and safety-by-construction are requirements. The trade-off, coverage, is measured honestly via the "
         "fallback rate - and the hybrid architecture is the planned comparison study.")

# ---------------------------------------------------------------- 4 · prototype
s = slide_new()
title(s, "The prototype in one view")
flow = [
    ("User message", TINT, LINE, INK),
    ("Safety layer\n(always first)", RED_TINT, RED, INK),
    ("Intent routing\n(deterministic)", TINT, LINE, INK),
    ("Knowledge base\n79 foods · 14 topics", TINT, LINE, INK),
    ("Grounded answer\n+ intent & source", TEAL, TEAL, WHITE),
]
x = 0.5
for i, (label, fill, ln, txtcol) in enumerate(flow):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.45, 1.62, 1.05, fill=fill, line=ln, radius=0.10)
    parts = label.split("\n")
    tb(s, x + 0.08, 1.62, 1.46, 0.8,
       [{"text": parts[0], "size": 11, "bold": True, "color": txtcol, "align": PP_ALIGN.CENTER, "after": 2}]
       + [{"text": p, "size": 9.5, "color": txtcol, "align": PP_ALIGN.CENTER} for p in parts[1:]])
    if i < 4:
        shape(s, MSO_SHAPE.RIGHT_ARROW, x + 1.66, 1.86, 0.22, 0.22, fill=TEAL)
    x += 1.84
tiles = [
    "Runs fully offline - zero external calls, local-only usage log",
    "Traffic-light rating (green / yellow / red) for everyday choices",
    "Bilingual (EN/DE) - ships as a windowed desktop app (.exe, no Python needed)",
]
for i, t in enumerate(tiles):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5 + i * 3.05, 3.0, 2.9, 1.15, fill=TINT, line=LINE, radius=0.08)
    tb(s, 0.72 + i * 3.05, 3.18, 2.46, 0.85, [{"text": t, "size": 11, "color": INK}])
tb(s, 0.5, 4.5, 9.0, 0.4, [{"text": "Every reply displays its detected intent and its source - explainability is built into the interface.",
                            "size": 11.5, "italic": True, "color": MUTED}])
footer(s, 4)
notes(s, "Walk the pipeline left to right and stress that the safety layer sits before everything else. "
         "The knowledge base is plain JSON - reviewable and extendable. Demo tip: run NutriCoach.exe live, "
         "ask 'white rice or brown rice?'.")

# ---------------------------------------------------------------- 5 · safety
s = slide_new()
title(s, "Safety by construction")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 1.4, 3.95, 0.7, fill=TEAL, radius=0.15)
tb(s, 0.72, 1.55, 3.5, 0.45, [{"text": "“How much insulin should I take with dinner?”",
                               "size": 11.5, "italic": True, "color": WHITE}])
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 2.3, 4.0, 1.75, fill=TINT, line=LINE, radius=0.07)
tb(s, 0.92, 2.48, 3.55, 1.45, [{"text": "“I'm deliberately not able to help with medication decisions - doses "
                                        "and timing belong with your doctor or diabetes team. If it's urgent, "
                                        "please contact your care team now.”", "size": 11, "color": INK}])
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 4.2, 2.1, 0.32, fill=WHITE, line=LINE, radius=0.5)
tb(s, 0.7, 4.26, 2.1, 0.22, [{"text": "intent: safety_medication", "size": 8.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
tb(s, 5.1, 1.45, 4.4, 3.3, [{"text": "A guarantee, not a guideline", "size": 15, "bold": True, "after": 10}] + bullets([
    "Safety patterns are checked before all other routing",
    "Three classes: medication · acute symptoms · crisis",
    "Answered only with a referral to human care - never with content",
    "Persistent \"not medical advice\" disclaimer on every reply",
    "Behaviour verified by unit checks and the evaluation set",
], size=12, after=8))
footer(s, 5)
notes(s, "Read the example: even though the question mentions dinner, it never reaches the meal-ideas route - "
         "the safety layer wins by priority. This is enforceable precisely because routing is code, not a prompt.")

# ---------------------------------------------------------------- 6 · evaluation
s = slide_new()
title(s, "Does the routing work? Measured.")
stats = [
    ("40 / 40", "unit checks passed", "routing in EN and DE, aliases, safety priority, knowledge-base integrity"),
    ("100 %", "routing accuracy", "78-utterance evaluation set across all intent families (78/78)"),
    ("11 / 11", "safety utterances deflected", "medication, acute-symptom and crisis inputs in the set"),
]
for i, (big, label, sub) in enumerate(stats):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5 + i * 3.1, 1.5, 2.9, 2.15, fill=TINT, line=LINE, radius=0.07)
    tb(s, 0.6 + i * 3.1, 1.78, 2.7, 1.7,
       [{"text": big, "size": 40, "bold": True, "color": TEAL, "font": "Cambria", "align": PP_ALIGN.CENTER, "after": 2},
        {"text": label, "size": 12.5, "bold": True, "align": PP_ALIGN.CENTER, "after": 6},
        {"text": sub, "size": 10, "color": MUTED, "align": PP_ALIGN.CENTER}])
tb(s, 0.5, 4.05, 9.0, 0.9, [
    {"text": "Honest caveat: the evaluation set is developer-authored - an upper bound, not field performance.", "size": 12, "bold": True, "after": 4},
    {"text": "Real-world coverage is therefore tracked separately: the fallback rate from the local usage log, plus every unanswered utterance as a concrete knowledge-base gap.", "size": 11.5, "color": MUTED}])
footer(s, 6)
notes(s, "Two evaluation layers: 27 unit checks and a 78-utterance routing evaluation, currently 100 percent. "
         "Be upfront that the set is developer-authored - the professor will respect the caveat. The field metric "
         "is the fallback rate computed from real usage logs.")

# ---------------------------------------------------------------- 7 · study design
s = slide_new()
title(s, "Evaluating impact: pilot RCT design")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 2.30, 2.0, 1.05, fill=TINT, line=LINE, radius=0.10)
tb(s, 0.6, 2.47, 1.8, 0.75, [{"text": "N = 120", "size": 15, "bold": True, "align": PP_ALIGN.CENTER, "after": 2},
                             {"text": "adults with T2D", "size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER}])
tb(s, 2.28, 2.02, 0.95, 0.3, [{"text": "randomized 1:1", "size": 9, "color": MUTED, "align": PP_ALIGN.CENTER}])
shape(s, MSO_SHAPE.RIGHT_ARROW, 2.56, 1.78, 0.5, 0.2, fill=ORANGE)
shape(s, MSO_SHAPE.RIGHT_ARROW, 2.56, 3.68, 0.5, 0.2, fill=BLUE)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.18, 1.42, 3.0, 0.95, fill=ORANGE_TINT, line=ORANGE, radius=0.10)
tb(s, 3.32, 1.56, 2.72, 0.7, [{"text": "Intervention · n = 60", "size": 12.5, "bold": True, "after": 2},
                              {"text": "usual care + NutriCoach chatbot", "size": 10.5, "color": MUTED}])
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.18, 3.30, 3.0, 0.95, fill=BLUE_TINT, line=BLUE, radius=0.10)
tb(s, 3.32, 3.44, 2.72, 0.7, [{"text": "Control · n = 60", "size": 12.5, "bold": True, "after": 2},
                              {"text": "usual care only", "size": 10.5, "color": MUTED}])
tb(s, 6.20, 2.02, 0.85, 0.3, [{"text": "12 weeks", "size": 9, "color": MUTED, "align": PP_ALIGN.CENTER}])
shape(s, MSO_SHAPE.RIGHT_ARROW, 6.30, 1.78, 0.5, 0.2, fill=ORANGE)
shape(s, MSO_SHAPE.RIGHT_ARROW, 6.30, 3.68, 0.5, 0.2, fill=BLUE)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.95, 1.62, 2.55, 2.45, fill=TINT, line=LINE, radius=0.08)
tb(s, 7.15, 1.82, 2.2, 2.1, [{"text": "Outcomes", "size": 13, "bold": True, "after": 8}] + bullets([
    "Primary: HbA1c change (week 0 to 12)",
    "Secondary: weekly fasting glucose",
    "Process: engagement (messages sent)",
], size=10.5, after=7))
tb(s, 0.5, 4.62, 9.0, 0.4, [{"text": "Demonstrated end-to-end on simulated data (documented assumptions, fixed seed) - the pipeline is pre-specified and reproducible.",
                             "size": 11.5, "italic": True, "color": MUTED}])
footer(s, 7)
notes(s, "Standard two-arm pilot design. The colors match the charts on the next slides: orange is intervention, "
         "blue is control. Emphasize pre-specification: the analysis plan exists as code before any real data.")

# ---------------------------------------------------------------- 8 · primary result
s = slide_new()
title(s, "Primary outcome (simulated pilot)", width=5.8)
badge_simulated(s)
tb(s, 0.5, 1.45, 4.3, 3.4, [
    {"text": "−0.37 pp", "size": 48, "bold": True, "color": TEAL, "font": "Cambria", "after": 2},
    {"text": "greater HbA1c reduction than control", "size": 13, "bold": True, "after": 12},
    {"text": "95% CI  [−0.59, −0.16]  ·  p < 0.001", "size": 12.5, "after": 6},
    {"text": "Effect size: Cohen's d = −0.63 (moderate)", "size": 12.5, "after": 12},
    {"text": "A 0.3-0.5 pp reduction is the range reported for real digital diabetes interventions - and is considered clinically relevant.",
     "size": 11, "color": MUTED}])
s.shapes.add_picture(os.path.join(FIG, "fig1_hba1c_change.png"), Inches(5.0), Inches(1.25), width=Inches(4.5))
footer(s, 8)
notes(s, "The intervention arm improves 0.37 percentage points more than control, Welch t-test, clearly significant, "
         "moderate effect size. Say it out loud every time: these numbers are simulated with assumed effects - "
         "they demonstrate the analysis, not efficacy.")

# ---------------------------------------------------------------- 9 · dose-response
s = slide_new()
title(s, "Engagement dose-response (simulated)", width=5.8)
badge_simulated(s)
tb(s, 0.5, 1.45, 4.3, 3.5, [
    {"text": "−0.064 pp", "size": 34, "bold": True, "color": TEAL, "font": "Cambria", "after": 2},
    {"text": "per 10 chatbot messages sent", "size": 13, "bold": True, "after": 12},
    {"text": "r = −0.36  ·  p = 0.005 (intervention arm)", "size": 12.5, "after": 12},
    {"text": "More engaged participants improve more - by construction of the simulation.", "size": 11.5, "after": 6},
    {"text": "In a real trial this association is observational: engaged users may differ. It supports, but cannot prove, causality.",
     "size": 11, "color": MUTED}])
s.shapes.add_picture(os.path.join(FIG, "fig3_engagement.png"), Inches(5.0), Inches(1.35), width=Inches(4.5))
footer(s, 9)
notes(s, "The secondary hypothesis: a dose-response between usage and improvement. Point out the honest caveat - "
         "even with real data this is observational; a causal claim needs a design-level answer.")

# ---------------------------------------------------------------- 10 · power
s = slide_new()
title(s, "Sizing the definitive trial")
stats = [
    ("63", "per arm", "to detect 0.3 pp (α = 0.05, power 80%, SD 0.6)"),
    ("79", "per arm", "with a 20% dropout reserve"),
    ("78 %", "pilot power", "the simulated pilot (n = 60/arm) for the same effect"),
]
for i, (big, label, sub) in enumerate(stats):
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5 + i * 3.1, 1.5, 2.9, 2.15, fill=TINT, line=LINE, radius=0.07)
    tb(s, 0.6 + i * 3.1, 1.74, 2.7, 1.75,
       [{"text": big, "size": 44, "bold": True, "color": TEAL, "font": "Cambria", "align": PP_ALIGN.CENTER, "after": 2},
        {"text": label, "size": 12.5, "bold": True, "align": PP_ALIGN.CENTER, "after": 6},
        {"text": sub, "size": 10, "color": MUTED, "align": PP_ALIGN.CENTER}])
tb(s, 0.5, 4.05, 9.0, 0.8, [
    {"text": "Normal-approximation two-sample sizing, computed in the analysis pipeline (analyze.py) - reproducible, and easy to re-run for other assumptions.",
     "size": 11.5, "color": MUTED}])
footer(s, 10)
notes(s, "The answer to the inevitable question 'how many participants would the real study need': about 63 per arm "
         "for the literature-typical 0.3 percentage points, 79 with dropout reserve. The pilot itself already has "
         "78 percent power for that effect under the planning assumptions.")

# ---------------------------------------------------------------- 11 · limitations & ethics
s = slide_new()
title(s, "Honest limitations & ethics")
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.5, 1.3, 4.4, 3.5, fill=TINT, line=LINE, radius=0.06)
tb(s, 0.85, 1.55, 3.75, 3.1, [{"text": "Limitations", "size": 15, "bold": True, "after": 10}] + bullets([
    "Study data are simulated - the pipeline, not the effect, is the contribution",
    "Coverage bounded by the knowledge base (fallback rate tracked)",
    "Nutrition values are curated approximations",
    "Population-level education only - no personalization (by design)",
], size=11.5, after=8))
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 5.1, 1.3, 4.4, 3.5, fill=TINT, line=LINE, radius=0.06)
tb(s, 5.45, 1.55, 3.75, 3.1, [{"text": "Ethics & privacy", "size": 15, "bold": True, "after": 10}] + bullets([
    "No real patient data anywhere in the project",
    "Not a medical device - persistent disclaimer, refuses medication/emergency questions",
    "Fully local processing; log stays on the device",
    "Real study path: ethics approval, informed consent, clinician review of all content",
], size=11.5, after=8))
footer(s, 11)
notes(s, "Owning the limitations is a strength in a research context. The ethics column doubles as the checklist "
         "for the moment real users enter the picture.")

# ---------------------------------------------------------------- 12 · roadmap
s = slide_new(TEAL_DARK)
tb(s, 0.5, 0.42, 9.0, 0.62, [{"text": "Where this goes next", "size": 30, "bold": True, "color": WHITE, "font": "Cambria"}])
steps = [
    "Real-user pilot: fallback rate, satisfaction (SUS), dietitian-rated guideline concordance",
    "Hybrid architecture: LLM for language, curated knowledge for content, same safety layer - comparison study",
    "Maintained food database (USDA FoodData Central), more languages beyond EN/DE",
    "The randomized pilot trial designed in this project",
]
for i, t in enumerate(steps):
    shape(s, MSO_SHAPE.OVAL, 0.6, 1.32 + i * 0.82, 0.44, 0.44, fill=WHITE)
    tb(s, 0.6, 1.40 + i * 0.82, 0.44, 0.3, [{"text": str(i + 1), "size": 15, "bold": True, "color": TEAL_DARK, "align": PP_ALIGN.CENTER}])
    tb(s, 1.25, 1.40 + i * 0.82, 8.2, 0.62, [{"text": t, "size": 13.5, "color": WHITE}])
tb(s, 0.5, 4.78, 9.0, 0.5, [{"text": "Every number in this deck is reproducible from the repository - code, data and analysis included.",
                             "size": 12, "italic": True, "color": TEAL_PALE}])
footer(s, 12, dark=True)
notes(s, "Close with the reproducibility line - it is the strongest sentence of the talk. Offer a live demo of the "
         "chatbot right after questions.")

prs.core_properties.title = "NutriCoach T2D - proof of concept"
prs.core_properties.author = "Sana Sajid"
prs.save(OUT)
print(f"Wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")
